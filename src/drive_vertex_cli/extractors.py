from __future__ import annotations

from html.parser import HTMLParser
from io import BytesIO
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


class UnsupportedFileTypeError(RuntimeError):
    """Raised when a downloaded Drive file has no configured text extractor."""

    pass


class HTMLTextParser(HTMLParser):
    """Very small HTML-to-text parser for lightweight document extraction."""

    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"br", "p", "div", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if data.strip():
            self.parts.append(data)

    def get_text(self) -> str:
        lines = [line.strip() for line in "".join(self.parts).splitlines()]
        return "\n".join(line for line in lines if line)


TEXT_MIME_PREFIXES = ("text/",)
TEXT_MIME_TYPES = {
    "application/json",
    "application/ld+json",
    "application/xml",
    "application/x-yaml",
    "application/yaml",
    "application/javascript",
    "application/x-sh",
}


def extract_text(name: str, mime_type: str, payload: bytes) -> str:
    """Convert a downloaded document payload into plain text."""

    suffix = Path(name).suffix.lower()

    if mime_type == "text/html" or suffix in {".html", ".htm"}:
        parser = HTMLTextParser()
        parser.feed(payload.decode("utf-8", errors="ignore"))
        return parser.get_text()

    if mime_type.startswith(TEXT_MIME_PREFIXES) or mime_type in TEXT_MIME_TYPES:
        return payload.decode("utf-8", errors="ignore")

    if mime_type == "application/pdf" or suffix == ".pdf":
        return _extract_pdf(payload)

    if (
        mime_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or suffix == ".docx"
    ):
        return _extract_docx(payload)

    if (
        mime_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or suffix == ".pptx"
    ):
        return _extract_pptx(payload)

    if (
        mime_type
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or suffix == ".xlsx"
    ):
        return _extract_xlsx(payload)

    raise UnsupportedFileTypeError(
        f"Unsupported file type for {name}: {mime_type or 'unknown'}"
    )


def _extract_pdf(payload: bytes) -> str:
    """Extract concatenated text from a PDF payload."""

    reader = PdfReader(BytesIO(payload))
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_docx(payload: bytes) -> str:
    """Extract paragraph text from a DOCX payload."""

    document = Document(BytesIO(payload))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs]
    return "\n".join(paragraph for paragraph in paragraphs if paragraph)


def _extract_pptx(payload: bytes) -> str:
    """Extract visible slide text from a PPTX payload."""

    presentation = Presentation(BytesIO(payload))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
        slide_text = "\n".join(parts).strip()
        if slide_text:
            slides.append(slide_text)
    return "\n\n".join(slides)


def _extract_xlsx(payload: bytes) -> str:
    """Extract a row-oriented text view from an XLSX payload."""

    workbook = load_workbook(BytesIO(payload), read_only=True, data_only=True)
    try:
        sheets: list[str] = []
        for worksheet in workbook.worksheets:
            rows: list[str] = [f"Sheet: {worksheet.title}"]
            for row in worksheet.iter_rows(values_only=True):
                values = [
                    str(value).strip()
                    for value in row
                    if value is not None and str(value).strip()
                ]
                if values:
                    rows.append(" | ".join(values))
            if len(rows) > 1:
                sheets.append("\n".join(rows))
        return "\n\n".join(sheets)
    finally:
        workbook.close()
